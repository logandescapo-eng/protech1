"""Generate the ProTech academic defence presentation (PPTX)."""

import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "documentattion" / "ProTech_Presentation.pptx"
CHART_DIR = ROOT / "documentattion" / "report_charts"

sys.path.insert(0, str(ROOT / "scripts"))
from growth_charts import generate_all_charts  # noqa: E402
from diagrams import generate_all_diagrams  # noqa: E402

# Brand palette
BLUE_DARK = RGBColor(0x1E, 0x3A, 0x8A)
BLUE_MID = RGBColor(0x25, 0x63, 0xEB)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_PALE = RGBColor(0xEF, 0xF6, 0xFF)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1E, 0x29, 0x3B)

AUTHOR = "Manye Patrice"
PROGRAMME = "IT Innovation in Business — Master's Degree"
PROJECT_TYPE = "End-of-Year Master's Project"
COURSES = "Web Technology in Business · Web Application Development"
REPO = "github.com/logandescapo-eng/protech1"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.55)
MARGIN_R = Inches(0.55)
CONTENT_TOP = Inches(1.35)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

_slide_num = 0


def _next_num():
    global _slide_num
    _slide_num += 1
    return _slide_num


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _textbox(slide, left, top, width, height, text, size=18, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _bullets(slide, left, top, width, height, items, size=16, color=BLACK, spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = spacing
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
    return box


def _footer(slide, num, total, backup=False):
    if backup:
        text = f"{AUTHOR}  |  BACKUP SLIDE  |  Slide {num}"
    else:
        text = f"{AUTHOR}  |  {PROGRAMME}  |  Slide {num} of {total}"
    _textbox(slide, MARGIN_L, Inches(7.05), CONTENT_W, Inches(0.35),
             text, size=9, color=SLATE, align=PP_ALIGN.CENTER)


def _header_bar(slide, title, subtitle=None):
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), BLUE_DARK)
    _textbox(slide, MARGIN_L, Inches(0.22), CONTENT_W, Inches(0.55),
             title, size=28, bold=True, color=WHITE)
    if subtitle:
        _textbox(slide, MARGIN_L, Inches(0.72), CONTENT_W, Inches(0.35),
                 subtitle, size=12, color=BLUE_LIGHT)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title_slide(prs, total):
    slide = _blank_slide(prs)
    _fill_bg(slide, BLUE_DARK)
    num = _next_num()

    _textbox(slide, MARGIN_L, Inches(1.6), CONTENT_W, Inches(1.0),
             "ProTech", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(2.55), CONTENT_W, Inches(0.55),
             "Professional Services Marketplace Platform", size=20, color=BLUE_LIGHT,
             align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(3.35), CONTENT_W, Inches(0.45),
             PROJECT_TYPE, size=16, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(3.85), CONTENT_W, Inches(0.4),
             PROGRAMME, size=14, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(4.3), CONTENT_W, Inches(0.35),
             COURSES, size=12, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    _add_rect(slide, Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.03), BLUE_MID)
    _textbox(slide, MARGIN_L, Inches(5.2), CONTENT_W, Inches(0.4),
             AUTHOR, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(5.65), CONTENT_W, Inches(0.35),
             date.today().strftime("%B %d, %Y"), size=13, color=BLUE_LIGHT,
             align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(6.2), CONTENT_W, Inches(0.3),
             REPO, size=10, color=SLATE, align=PP_ALIGN.CENTER)

    _notes(slide,
           "Welcome. Introduce yourself, the programme, and the project title. "
           "Mention this is an 8–10 minute defence of a full-stack marketplace platform.")
    return num


def section_slide(prs, title, subtitle, total, backup=False):
    slide = _blank_slide(prs)
    _fill_bg(slide, INDIGO)
    num = _next_num()
    _textbox(slide, MARGIN_L, Inches(2.8), CONTENT_W, Inches(0.9),
             title, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(3.75), CONTENT_W, Inches(0.5),
             subtitle, size=16, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    _footer(slide, num, total, backup=backup)
    return num


def backup_divider_slide(prs, title, total):
    slide = _blank_slide(prs)
    _fill_bg(slide, SLATE)
    num = _next_num()
    _textbox(slide, MARGIN_L, Inches(2.9), CONTENT_W, Inches(0.8),
             "BACKUP SLIDES", size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(3.75), CONTENT_W, Inches(0.5),
             title, size=14, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(4.4), CONTENT_W, Inches(0.4),
             "Hidden from main slideshow — use during Q&A if needed",
             size=11, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    _footer(slide, num, total, backup=True)
    return num


def hide_slides(prs, indices):
    """Hide slide indices (0-based) from slideshow (show='0' on sldId)."""
    sld_ids = prs.slides._sldIdLst
    for i in indices:
        if 0 <= i < len(sld_ids):
            sld_ids[i].set("show", "0")


def content_slide(prs, title, bullets, total, notes="", subtitle=None,
                  bullet_size=16, backup=False):
    slide = _blank_slide(prs)
    num = _next_num()
    _header_bar(slide, title, subtitle)
    _bullets(slide, MARGIN_L, CONTENT_TOP, CONTENT_W, Inches(5.5),
             bullets, size=bullet_size)
    _footer(slide, num, total, backup=backup)
    if notes:
        _notes(slide, notes)
    return num


def two_column_slide(prs, title, left_items, right_items, total, notes="",
                     left_title="", right_title="", backup=False):
    slide = _blank_slide(prs)
    num = _next_num()
    _header_bar(slide, title)
    col_w = Inches(5.85)
    gap = Inches(0.35)
    right_left = MARGIN_L + col_w + gap

    if left_title:
        _textbox(slide, MARGIN_L, CONTENT_TOP, col_w, Inches(0.35),
                 left_title, size=14, bold=True, color=BLUE_DARK)
    if right_title:
        _textbox(slide, right_left, CONTENT_TOP, col_w, Inches(0.35),
                 right_title, size=14, bold=True, color=BLUE_DARK)

    top = CONTENT_TOP + Inches(0.4)
    _bullets(slide, MARGIN_L, top, col_w, Inches(5.0), left_items, size=15)
    _bullets(slide, right_left, top, col_w, Inches(5.0), right_items, size=15)
    _footer(slide, num, total, backup=backup)
    if notes:
        _notes(slide, notes)
    return num


def table_slide(prs, title, headers, rows, total, notes="", col_widths=None, backup=False):
    slide = _blank_slide(prs)
    num = _next_num()
    _header_bar(slide, title)
    cols = len(headers)
    tbl_rows = len(rows) + 1
    left = MARGIN_L
    top = CONTENT_TOP
    width = CONTENT_W
    height = Inches(0.45 * tbl_rows)

    table = slide.shapes.add_table(tbl_rows, cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(rows):
        fill = BLUE_PALE if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = BLACK

    _footer(slide, num, total, backup=backup)
    if notes:
        _notes(slide, notes)
    return num


def chart_slide(prs, title, image_path, caption, total, notes="", backup=False):
    slide = _blank_slide(prs)
    num = _next_num()
    _header_bar(slide, title)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), MARGIN_L, CONTENT_TOP, width=CONTENT_W)
    else:
        _add_rect(slide, MARGIN_L, CONTENT_TOP, CONTENT_W, Inches(4.8),
                  BLUE_PALE, BLUE_MID, Pt(2))
        _textbox(slide, MARGIN_L, Inches(3.2), CONTENT_W, Inches(0.5),
                 "[ Chart not found — run scripts/growth_charts.py ]",
                 size=14, color=SLATE, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(6.35), CONTENT_W, Inches(0.35),
             caption, size=10, color=SLATE, align=PP_ALIGN.CENTER)
    _footer(slide, num, total, backup=backup)
    if notes:
        _notes(slide, notes)
    return num


def screenshot_slide(prs, fig_num, title, url, login, total, notes="", backup=False):
    slide = _blank_slide(prs)
    num = _next_num()
    _header_bar(slide, f"Figure {fig_num} — {title}")

    frame_top = CONTENT_TOP
    frame_h = Inches(4.55)
    shape = _add_rect(slide, MARGIN_L, frame_top, CONTENT_W, frame_h,
                      RGBColor(0xF8, 0xFA, 0xFC), BLUE_MID, Pt(1.5))
    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    _textbox(slide, MARGIN_L, Inches(2.85), CONTENT_W, Inches(0.45),
             "[ INSERT SCREENSHOT HERE ]", size=18, bold=True, color=BLUE_MID,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    meta = f"URL: {url}"
    if login:
        meta += f"   |   Login: {login}"
    _textbox(slide, MARGIN_L, Inches(5.95), CONTENT_W, Inches(0.3),
             meta, size=9, color=BLUE_DARK, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(6.3), CONTENT_W, Inches(0.3),
             "PowerPoint: Insert → Pictures, or drag image into the dashed frame",
             size=9, color=SLATE, align=PP_ALIGN.CENTER)
    _footer(slide, num, total, backup=backup)
    if notes:
        _notes(slide, notes)
    return num


def architecture_slide(prs, total, backup=False):
    slide = _blank_slide(prs)
    num = _next_num()
    _header_bar(slide, "System Architecture", "Three-tier design — Docker Compose deployment")

    layers = [
        (BLUE_MID, "Presentation", "Browser → Nginx (port 80) → static files & reverse proxy"),
        (INDIGO, "Application", "Gunicorn → Django 5.2  |  users · bookings · reviews · escrow"),
        (RGBColor(0x0E, 0xA5, 0xE9), "Data & Cache", "PostgreSQL 15  +  Redis 7 (sessions & cache)"),
        (SLATE, "Cross-cutting", "RBAC · Escrow ledger · Logging · Django Admin"),
    ]
    top = CONTENT_TOP
    box_h = Inches(0.95)
    gap = Inches(0.18)
    arrow_h = Inches(0.22)

    for color, name, desc in layers:
        _add_rect(slide, MARGIN_L + Inches(1.2), top, Inches(10.0), box_h, color)
        _textbox(slide, MARGIN_L + Inches(1.35), top + Inches(0.1), Inches(2.5), Inches(0.35),
                 name, size=14, bold=True, color=WHITE)
        _textbox(slide, MARGIN_L + Inches(3.9), top + Inches(0.18), Inches(7.0), Inches(0.6),
                 desc, size=12, color=WHITE)
        top += box_h + gap
        if top < Inches(6.0):
            _textbox(slide, MARGIN_L + Inches(6.0), top - gap + Inches(0.02),
                     Inches(0.5), Inches(0.2), "▼", size=14, bold=True, color=BLUE_MID,
                     align=PP_ALIGN.CENTER)
            top += arrow_h

    _footer(slide, num, total, backup=backup)
    _notes(slide,
           "Walk through each tier. Emphasise Nginx offloads static files, Django handles business logic, "
           "PostgreSQL is source of truth, Redis speeds up sessions and worker search.")
    return num


def image_slide(prs, title, image_path, caption, total, notes="", backup=False):
    """Full-width diagram slide for ER / sequence diagrams."""
    slide = _blank_slide(prs)
    num = _next_num()
    _header_bar(slide, title)
    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), MARGIN_L, CONTENT_TOP, width=CONTENT_W)
    _textbox(slide, MARGIN_L, Inches(6.35), CONTENT_W, Inches(0.35),
             caption, size=9, color=SLATE, align=PP_ALIGN.CENTER)
    _footer(slide, num, total, backup=backup)
    if notes:
        _notes(slide, notes)
    return num


def build_presentation():
    print("Generating growth charts and diagrams...")
    charts = generate_all_charts()
    diagrams = generate_all_diagrams()

    MAIN_TOTAL = 28  # slides in the main defence deck (excludes hidden backup)

    global _slide_num
    _slide_num = 0

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── MAIN DECK (28 slides) ─────────────────────────────────────────
    title_slide(prs, MAIN_TOTAL)

    content_slide(prs, "Presentation Outline", [
        "1.  Introduction, problem & goals",
        "2.  Business model & growth projections",
        "3.  System design & architecture",
        "4.  Implementation & security",
        "5.  Live demo (screenshots)",
        "6.  Testing, deployment & rubric alignment",
        "7.  Conclusion & Q&A",
    ], MAIN_TOTAL, notes="~8–10 minute defence. Backup slides follow for Q&A.")

    section_slide(prs, "1", "Introduction & Business Case", MAIN_TOTAL)

    content_slide(prs, "Introduction & Problem Statement", [
        "ProTech — a marketplace connecting clients with skilled service workers.",
        "End-of-year master's project: IT Innovation in Business.",
        "Problems: poor discovery, payment risk, fragmented coordination, no admin oversight.",
        "Solution: one Django platform with booking, escrow, reviews, and RBAC.",
        "Clone-and-run via Docker for academic evaluation.",
    ], MAIN_TOTAL, notes="Keep to 60 seconds. Make the problem relatable.")

    two_column_slide(prs, "Goals & Business Use Case",
                     [
                         "Project goals:",
                         "• Working client + worker marketplace",
                         "• Django, Docker, Redis, tests, logging",
                         "• Escrow with platform commission",
                         "• Django Admin CMS",
                     ],
                     [
                         "Who benefits:",
                         "• Clients — search, book, pay safely",
                         "• Workers — profile, jobs, wallet payouts",
                         "• Platform — commission + subscriptions",
                     ],
                     MAIN_TOTAL,
                     left_title="Technical goals", right_title="Stakeholders")

    section_slide(prs, "2", "Business Model & Growth", MAIN_TOTAL)

    content_slide(prs, "Business Model Overview", [
        "Hybrid revenue: transaction commission + worker subscriptions.",
        "Clients pay no membership — job price held in escrow until completion.",
        "Commission implemented in escrow/services.py (default 5%).",
        "Subscriptions (Starter / Professional / Premium) specified for production.",
        "Charts below are illustrative projections, not audited forecasts.",
    ], MAIN_TOTAL)

    two_column_slide(prs, "Revenue Model",
                     [
                         "Commission (on job completion):",
                         "• Micro (<$50): 8% or min $3",
                         "• Standard: 5%",
                         "• Commercial (>$500): 3.5%",
                     ],
                     [
                         "Worker subscriptions:",
                         "• Starter — free, basic listing",
                         "• Professional — $29/mo, −1% commission",
                         "• Premium — $59/mo, featured placement",
                     ],
                     MAIN_TOTAL,
                     left_title="Per transaction", right_title="Monthly plans")

    chart_slide(prs, "Traffic Growth Projection", charts["traffic"],
                "Month 24: ~9,500 visitors/mo (illustrative)", MAIN_TOTAL)
    chart_slide(prs, "Revenue Growth (MRR)", charts["revenue"],
                "Month 24: ~$50k MRR — commission + subscriptions", MAIN_TOTAL)
    chart_slide(prs, "Bookings & GMV", charts["bookings"],
                "Completed jobs and gross merchandise volume", MAIN_TOTAL)

    section_slide(prs, "3", "System Design & Architecture", MAIN_TOTAL)

    architecture_slide(prs, MAIN_TOTAL)

    table_slide(prs, "Technology Stack",
                ["Layer", "Technology", "Role"],
                [
                    ("Back-end", "Django 5.2.15", "ORM, auth, admin, security"),
                    ("Front-end", "Jinja2 templates", "Server-rendered UI"),
                    ("Database", "PostgreSQL 15", "Relational data + escrow ledger"),
                    ("Cache", "Redis 7", "Sessions and worker search"),
                    ("Deploy", "Docker Compose + Nginx", "Four-service stack"),
                ],
                MAIN_TOTAL, col_widths=[1.5, 2.5, 4.5])

    content_slide(prs, "Escrow Payment Flow", [
        "1. Client deposits demo funds → UserWallet.",
        "2. Client books worker → Booking (pending).",
        "3. Client funds escrow → EscrowHold created.",
        "4. Worker completes job → release_escrow() runs.",
        "5. Worker paid net of platform fee; transactions logged.",
        "6. Cancel → refund_escrow() returns funds to client.",
    ], MAIN_TOTAL, bullet_size=15,
    notes="Key differentiator. Mention Decimal + atomic DB transactions.")

    section_slide(prs, "4", "Implementation & Security", MAIN_TOTAL)

    two_column_slide(prs, "Security & Implementation",
                     [
                         "Security:",
                         "• Email login, PBKDF2 hashing",
                         "• CSRF, ORM, XSS escaping",
                         "• RBAC (@user_type_required)",
                         "• Secrets in .env only",
                     ],
                     [
                         "Code structure:",
                         "• users/ bookings/ reviews/ escrow/",
                         "• escrow/services.py — all money logic",
                         "• Redis cache (300s worker search)",
                         "• Logging → stdout + django.log",
                     ],
                     MAIN_TOTAL,
                     left_title="Controls", right_title="Django apps")

    section_slide(prs, "5", "Application Demo", MAIN_TOTAL)

    screenshot_slide(prs, 1, "Landing Page (Home)", "http://localhost/",
                     None, MAIN_TOTAL)
    screenshot_slide(prs, 2, "Client Dashboard", "http://localhost/user/dashboard/",
                     "john@example.com / password123", MAIN_TOTAL)
    screenshot_slide(prs, 3, "Worker Dashboard", "http://localhost/worker/dashboard/",
                     "mike@example.com / password123", MAIN_TOTAL)
    screenshot_slide(prs, 4, "Escrow Payment", "http://localhost/booking/<id>/escrow/",
                     "john@example.com / password123", MAIN_TOTAL)

    section_slide(prs, "6", "Testing & Deployment", MAIN_TOTAL)

    content_slide(prs, "Testing, Deployment & Demo", [
        "17 automated tests pass (python manage.py test).",
        "GitHub Actions: flake8, migrate, test, /health/ on every push.",
        "Docker: db · redis · backend · nginx — docker compose up --build.",
        "Demo: john@example.com / mike@example.com — password123.",
        "Health: http://localhost/health/ → {\"status\": \"ok\"}",
    ], MAIN_TOTAL, bullet_size=15)

    table_slide(prs, "Rubric Alignment (Summary)",
                ["Requirement", "Status", "Evidence"],
                [
                    ("Django + templates", "✓", "templates/, Jinja2 env"),
                    ("Auth & RBAC", "✓", "users/decorators.py"),
                    ("Docker + Redis", "✓", "docker-compose.yml"),
                    ("Tests + CI", "✓", "17 tests, ci.yml"),
                    ("Logging + security", "✓", "settings.py, CSRF"),
                ],
                MAIN_TOTAL, col_widths=[2.0, 0.8, 3.7],
                notes="Full rubric table is in the technical report Appendix D.")

    section_slide(prs, "7", "Conclusion", MAIN_TOTAL)

    two_column_slide(prs, "Lessons, Future Work & Conclusion",
                     [
                         "Lessons learned:",
                         "• Domain-split Django apps scale well",
                         "• Flask→Django needed Jinja2 shim",
                         "• Escrow needs Decimal + TXNs",
                         "• Docker entrypoint seeds demo data",
                     ],
                     [
                         "Conclusion:",
                         "• Full-stack brief met",
                         "• Working escrow marketplace",
                         "• Clear business model + growth path",
                         "• Ready for live demo & Q&A",
                     ],
                     MAIN_TOTAL,
                     left_title="Reflection", right_title="Outcome",
                     notes="Thank the panel. Offer live demo if time allows.")

    # Closing slide
    slide = _blank_slide(prs)
    num = _next_num()
    _fill_bg(slide, BLUE_DARK)
    _textbox(slide, MARGIN_L, Inches(2.8), CONTENT_W, Inches(0.9),
             "Thank You", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(3.85), CONTENT_W, Inches(0.5),
             "Questions & Discussion", size=20, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(4.8), CONTENT_W, Inches(0.4),
             f"{AUTHOR}  ·  {PROGRAMME}", size=14, color=WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, MARGIN_L, Inches(5.3), CONTENT_W, Inches(0.35),
             REPO, size=11, color=SLATE, align=PP_ALIGN.CENTER)
    _notes(slide, "Pause for questions. Have demo ready at http://localhost.")

    main_slide_count = num
    backup_start_index = main_slide_count  # 0-based index of first backup slide

    # ── BACKUP SLIDES (hidden from slideshow) ─────────────────────────
    backup_total = MAIN_TOTAL  # footer label only

    backup_divider_slide(prs, "Additional material for Q&A", backup_total)

    chart_slide(prs, "BACKUP — User Growth", charts["users"],
                "Registered client and worker accounts (24 months)", backup_total, backup=True)
    chart_slide(prs, "BACKUP — Revenue Mix", charts["mix"],
                "Subscription share of MRR after Month 7", backup_total, backup=True)

    table_slide(prs, "BACKUP — Full Use Case List",
                ["ID", "Use case", "Actor"],
                [
                    ("UC-01", "Register & login", "Visitor"),
                    ("UC-04", "Reset password", "User"),
                    ("UC-05", "Browse workers", "Client"),
                    ("UC-06", "Create booking", "Client"),
                    ("UC-07", "Manage booking status", "Client / Worker"),
                    ("UC-08", "Fund escrow", "Client"),
                    ("UC-09", "Release / refund escrow", "System"),
                    ("UC-10", "Leave review", "Client"),
                    ("UC-11", "Admin CRUD", "Staff"),
                ],
                backup_total, col_widths=[0.7, 3.5, 1.3], backup=True)

    image_slide(prs, "BACKUP — ER Diagram", diagrams["er"],
                "Core database schema (users, bookings, escrow, reviews)", backup_total, backup=True)
    image_slide(prs, "BACKUP — Escrow Sequence", diagrams["sequence"],
                "Runtime flow: deposit → book → fund → complete → release", backup_total, backup=True)

    screenshot_slide(prs, 5, "Browse Workers", "http://localhost/workers/",
                     "john@example.com / password123", backup_total, backup=True)
    screenshot_slide(prs, 6, "Django Admin", "http://localhost/admin/",
                     "Superuser account", backup_total, backup=True)

    table_slide(prs, "BACKUP — Full Rubric Alignment",
                ["Requirement", "Implementation", "Evidence"],
                [
                    ("Django + templates", "Django 5.2 + Jinja2", "templates/"),
                    ("Auth & RBAC", "Login, @user_type_required", "users/"),
                    ("Password reset", "/password-reset/", "users/urls.py"),
                    ("Admin CMS", "Django Admin", "/admin/"),
                    ("Security", "CSRF, hashing, .env", "settings.py"),
                    ("Tests", "17 passing", "manage.py test"),
                    ("Docker", "4 services", "docker-compose.yml"),
                    ("Redis cache", "cache_utils.py", "300s worker cache"),
                    ("Logging", "Per-app loggers", "logs/django.log"),
                    ("Linting", "flake8 + black", "ci.yml"),
                ],
                backup_total, col_widths=[1.5, 2.0, 2.8], backup=True)

    content_slide(prs, "BACKUP — Limitations", [
        "Demo wallet only — no live payment provider integrated.",
        "Worker subscriptions specified but not yet coded.",
        "No WebSocket real-time messaging.",
        "17 tests — no browser E2E suite yet.",
        "Growth charts are illustrative, not market research.",
    ], backup_total, bullet_size=14, backup=True)

    content_slide(prs, "BACKUP — References", [
        "Rochet & Tirole (2003) — two-sided markets",
        "Evans & Schmalensee (2016) — Matchmakers",
        "OWASP Top Ten (2021) — web security",
        "Django, PostgreSQL, Redis, Docker documentation",
        "Fielding (2000); Fowler (2002) — architecture patterns",
    ], backup_total, bullet_size=13, backup=True)

    # Hide all backup slides (from backup divider onward)
    hide_indices = list(range(backup_start_index, len(prs.slides)))
    hide_slides(prs, hide_indices)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Presentation written to: {OUTPUT}")
    print(f"Main slides: {main_slide_count}  |  Backup slides: {len(hide_indices)} (hidden)")
    print(f"Total slides in file: {len(prs.slides)}")
    return OUTPUT


if __name__ == "__main__":
    build_presentation()
